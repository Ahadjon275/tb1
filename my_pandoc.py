```python
import io
import os
import subprocess
from PyPDF2 import PdfReader
import pandas as pd
from pptx import Presentation
import my_log  # Предполагается, что этот модуль доступен
import my_ocr  # Предполагается, что этот модуль доступен
import utils    # Предполагается, что этот модуль доступен

def fb2_to_text(data, ext, lang):
    if ext == 'epub':
        # Конвертация EPUB в текст с использованием pandoc
        result = subprocess.run(['pandoc', '-f', 'epub', '-t', 'plain'], input=data, capture_output=True)
        return result.stdout.decode('utf-8')

    elif ext == 'pptx':
        return read_pptx(data)

    elif ext in ['docx', 'odt', 'rtf']:
        # Конвертация DOCX/ODT/RTF в текст с использованием pandoc
        result = subprocess.run(['pandoc', '-f', ext, '-t', 'plain'], input=data, capture_output=True)
        return result.stdout.decode('utf-8')

    elif ext == 'djvu':
        pdf_file = convert_djvu2pdf(data)
        if pdf_file:
            text = my_ocr.get_text_from_pdf(pdf_file)
            os.remove(pdf_file)  # Удаление временного PDF файла
            return text

    elif ext == 'pdf':
        reader = PdfReader(io.BytesIO(data))
        text = ''
        for page in reader.pages:
            text += page.extract_text() + '\n'
        return text

    elif ext in ['xlsx', 'xls']:
        df = pd.read_excel(io.BytesIO(data))
        return df.to_csv(index=False)

    else:
        raise ValueError(f"Unsupported file extension: {ext}")

def read_pptx(data):
    presentation = Presentation(io.BytesIO(data))
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def convert_djvu2pdf(data):
    temp_djvu_path = "temp.djvu"
    temp_pdf_path = "temp.pdf"
    
    with open(temp_djvu_path, 'wb') as f:
        f.write(data)

    try:
        subprocess.run(['ddjvu', '-format=pdf', temp_djvu_path, temp_pdf_path], check=True)
        return temp_pdf_path
    except subprocess.CalledProcessError as e:
        my_log.log2('Failed to convert DJVU to PDF: {}'.format(e))
        return None
    finally:
        os.remove(temp_djvu_path)  # Удаление временного DJVU файла

if __name__ == '__main__':
    # Здесь можно добавить тесты или основной функционал
    pass
```
